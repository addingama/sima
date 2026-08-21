#!/usr/bin/env python3
"""Generate SIMA product demo PowerPoint.

  python3 -m venv .venv-pptx
  .venv-pptx/bin/pip install python-pptx
  .venv-pptx/bin/python docs/demo/generate_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "SIMA-Presentasi-Demo.pptx"

NAVY = RGBColor(0x0F, 0x2C, 0x4C)
TEAL = RGBColor(0x0D, 0x7A, 0x6F)
SLATE = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)


def set_run(run, size=18, bold=False, color=SLATE):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_bg(slide, color=WHITE):
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    fill.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = fill._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def add_footer(slide, page: int, total: int):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = f"SIMA — Sistem Informasi Manajemen Amanah  ·  Demo produk  ·  {page}/{total}"
    set_run(p.runs[0], size=11, color=MUTED)


def add_accent_bar(slide):
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.9), Inches(13.333), Inches(1.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    p.text = "SIMA"
    set_run(p.runs[0], size=54, bold=True, color=WHITE)

    s = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.5))
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sistem Informasi Manajemen Amanah"
    set_run(p.runs[0], size=28, color=WHITE)
    p2 = tf.add_paragraph()
    p2.text = (
        "Kelola dana titipan dengan amanah: setiap uang masuk punya tujuan,\n"
        "dan setiap pengeluaran mengambil dari dana yang sesuai."
    )
    set_run(p2.runs[0], size=16, color=RGBColor(0xCB, 0xD5, 0xE1))

    f = slide.shapes.add_textbox(Inches(0.8), Inches(6.25), Inches(11.5), Inches(0.8))
    p = f.text_frame.paragraphs[0]
    p.text = "Materi demonstrasi produk  ·  untuk pimpinan, bendahara, dan tim IT"
    set_run(p.runs[0], size=14, color=WHITE)


def section_slide(prs, title, subtitle, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_accent_bar(slide)

    t = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(11.5), Inches(1))
    p = t.text_frame.paragraphs[0]
    p.text = title
    set_run(p.runs[0], size=36, bold=True, color=NAVY)

    s = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(11.5), Inches(1.2))
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    set_run(p.runs[0], size=18, color=SLATE)
    add_footer(slide, page, total)


def bullets_slide(prs, title, bullets, page, total, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_accent_bar(slide)

    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.7))
    p = t.text_frame.paragraphs[0]
    p.text = title
    set_run(p.runs[0], size=28, bold=True, color=NAVY)

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11.8), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(10)
        set_run(p.runs[0], size=17, color=SLATE)

    if notes:
        n = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(11.8), Inches(0.4))
        p = n.text_frame.paragraphs[0]
        p.text = notes
        set_run(p.runs[0], size=12, color=MUTED)

    add_footer(slide, page, total)


def two_col_slide(prs, title, left_title, left_items, right_title, right_items, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_accent_bar(slide)

    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.6))
    p = t.text_frame.paragraphs[0]
    p.text = title
    set_run(p.runs[0], size=26, bold=True, color=NAVY)

    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.2))
    left.fill.solid()
    left.fill.fore_color.rgb = LIGHT
    left.line.fill.background()

    lt = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(5.2), Inches(0.5))
    p = lt.text_frame.paragraphs[0]
    p.text = left_title
    set_run(p.runs[0], size=18, bold=True, color=TEAL)

    lb = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(5.2), Inches(4.1))
    tf = lb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + line
        p.space_after = Pt(8)
        set_run(p.runs[0], size=15, color=SLATE)

    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.2))
    right.fill.solid()
    right.fill.fore_color.rgb = LIGHT
    right.line.fill.background()

    rt = slide.shapes.add_textbox(Inches(7.1), Inches(1.4), Inches(5.2), Inches(0.5))
    p = rt.text_frame.paragraphs[0]
    p.text = right_title
    set_run(p.runs[0], size=18, bold=True, color=TEAL)

    rb = slide.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.1))
    tf = rb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + line
        p.space_after = Pt(8)
        set_run(p.runs[0], size=15, color=SLATE)

    add_footer(slide, page, total)


def table_slide(prs, title, headers, rows, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_accent_bar(slide)

    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.6))
    p = t.text_frame.paragraphs[0]
    p.text = title
    set_run(p.runs[0], size=26, bold=True, color=NAVY)

    cols = len(headers)
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        cols,
        Inches(0.6),
        Inches(1.2),
        Inches(12.1),
        Inches(0.55 * (len(rows) + 1)),
    )
    table = table_shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                set_run(run, size=13, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    set_run(run, size=12, color=SLATE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else LIGHT

    add_footer(slide, page, total)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = 13

    title_slide(prs)

    bullets_slide(
        prs,
        "Agenda (15–20 menit)",
        [
            "1. Masalah lembaga sosial: dana titipan mudah “tercecer” antar tujuan",
            "2. Konsep SIMA: Kas/Bank vs Dana Amanah + Amanah Ledger",
            "3. Siapa memakai sistem (role) dan alur persetujuan",
            "4. Live demo: master → penerimaan → pengeluaran → laporan",
            "5. Fitur pendukung: transfer, biaya bank, rekonsiliasi, portal donatur",
            "6. FAQ & langkah go-live",
        ],
        2,
        total,
    )

    bullets_slide(
        prs,
        "Masalah yang diselesaikan",
        [
            "Donasi masuk dengan niat tertentu (zakat, yatim, umum) — mudah tercampur di Excel/kas tunggal",
            "Pengeluaran sulit dibuktikan dari dana mana; sulit jawab auditor/donatur",
            "Saldo “kas” vs “peruntukan” sering tidak cocok; koreksi dengan edit/hapus baris",
            "Persetujuan ketua/verifikator tidak terdokumentasi rapi",
            "SIMA memaksa: uang masuk harus beralokasi; uang keluar harus dari dana yang tepat",
        ],
        3,
        total,
    )

    two_col_slide(
        prs,
        "Dua dimensi uang (inti produk)",
        "Kas / Bank — di mana uang berada",
        [
            "Kas tunai, rekening BCA, dll.",
            "Saldo fisik yang bisa direkonsiliasi dengan rekening koran",
            "Transfer antar rekening tidak mengubah Dana Amanah",
        ],
        "Dana Amanah — untuk apa uang boleh dipakai",
        [
            "Restricted: terikat niat (zakat, yatim, …)",
            "Unrestricted: dana umum lembaga",
            "Pengeluaran & biaya bank harus ambil dari dana yang sesuai",
        ],
        4,
        total,
    )

    bullets_slide(
        prs,
        "Amanah Ledger = sumber kebenaran tunggal",
        [
            "Setiap transaksi finansial = jurnal double-entry (kas/bank ↔ dana)",
            "Saldo dihitung dari ledger — bukan angka yang diedit manual",
            "Tidak ada hard delete; koreksi hanya lewat void/reversal",
            "Total Kas/Bank harus = Total Dana Amanah (selisih = 0)",
            "Audit trail: siapa mengubah apa, kapan; approval tercatat",
        ],
        5,
        total,
        notes="Invariant: saldo akun & dana tidak boleh negatif.",
    )

    table_slide(
        prs,
        "Alur uang singkat",
        ["Transaksi", "Efek Kas/Bank", "Efek Dana Amanah"],
        [
            ["Penerimaan (disetujui)", "Bertambah", "Bertambah di dana tujuan (alokasi inline)"],
            ["Pengeluaran (disetujui)", "Berkurang", "Berkurang dari sumber dana"],
            ["Biaya bank (diposting)", "Berkurang", "Berkurang dari dana operasional (bukan restricted)"],
            ["Transfer antar rekening", "Pindah antar akun", "Tidak berubah"],
            ["Reversal / void", "Negasi transaksi sumber", "Negasi transaksi sumber"],
        ],
        6,
        total,
    )

    table_slide(
        prs,
        "Peran pengguna (RBAC)",
        ["Role", "Fokus di demo", "Contoh aksi"],
        [
            ["Admin", "Setup & go-live", "User, saldo awal, semua permission"],
            ["Bendahara", "Operasional harian", "Master, draft/submit transaksi, biaya bank, transfer"],
            ["Verifikator", "Cek pengeluaran", "Verify / reject pengeluaran"],
            ["Ketua", "Persetujuan akhir", "Approve penerimaan & pengeluaran"],
            ["Auditor", "Pengawasan", "Laporan & audit trail (view)"],
            ["Donatur", "Portal", "Lihat ringkasan & riwayat donasi sendiri"],
        ],
        7,
        total,
    )

    bullets_slide(
        prs,
        "Script live demo (urutan disarankan)",
        [
            "Login bendahara → Dashboard (metric + chart real)",
            "Master: Donatur, Dana Amanah, Kas/Bank, Vendor, Event (jelaskan restricted)",
            "Penerimaan: buat + alokasi → submit → login ketua → approve → saldo berubah",
            "Pengeluaran: dari dana yang tepat + vendor → submit → verifikator → ketua approve",
            "Tunjukkan Approval queue, Laporan saldo, Rekonsiliasi global di dashboard",
            "Bonus: Transfer, Biaya bank, Portal donatur@sima.test",
        ],
        8,
        total,
        notes="Akun demo: *@sima.test / password · isi data: php artisan sima:seed-demo",
    )

    bullets_slide(
        prs,
        "Modul yang siap didemokan",
        [
            "Master: Donatur, Vendor, Dana Amanah, Kas/Bank, Event/Program, Users",
            "Keuangan: Penerimaan, Pengeluaran, Biaya Bank, Transfer, Rekonsiliasi Bank",
            "Operasional: Liabilitas, Saldo Awal (go-live), Approval",
            "Laporan: Saldo akun/dana, buku besar, opening, dashboard analitik",
            "Portal Donatur: profil + riwayat donasi approved",
            "API + Swagger + Postman untuk integrasi / audit teknis",
        ],
        9,
        total,
    )

    bullets_slide(
        prs,
        "Keamanan & kepatuhan (poin jual)",
        [
            "RBAC permission granular (Spatie) — menu & API ikut permission",
            "Ledger immutable (guard model + trigger DB)",
            "Workflow approval multi-tahap untuk pengeluaran",
            "Idempotency & nomor dokumen unik (RCP/DSB/FEE/…)",
            "Audit log + riwayat approval untuk auditor",
            "Rekonsiliasi: bukti selisih kas↔dana langsung di dashboard",
        ],
        10,
        total,
    )

    bullets_slide(
        prs,
        "Go-live singkat (pesan ke pimpinan)",
        [
            "Fase 0: sepakati daftar dana, rekening, pemetaan saldo awal",
            "Fase 1–3: user & master data (admin + bendahara)",
            "Fase 4: posting saldo awal (admin) — cutover",
            "Fase 5: transaksi harian + approval ketua/verifikator",
            "Detail langkah: docs/PANDUAN-MULAI.md",
            "FAQ demo & pertanyaan sulit: docs/demo/FAQ-DEMO.md",
        ],
        11,
        total,
    )

    bullets_slide(
        prs,
        "Ringkasan satu kalimat",
        [
            "SIMA memastikan dana titipan tetap amanah:",
            "uang punya tempat (kas/bank) dan punya tujuan (Dana Amanah),",
            "semua perubahan lewat ledger yang bisa diaudit,",
            "dan koreksi tidak pernah dengan menghapus sejarah.",
        ],
        12,
        total,
    )

    section_slide(
        prs,
        "Terima kasih",
        "Siap tanya jawab. Buka FAQ-DEMO.md untuk jawaban siap pakai.",
        13,
        total,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
